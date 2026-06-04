from pptx import Presentation
from sentence_transformers import SentenceTransformer
from langchain_text_splitters import RecursiveCharacterTextSplitter
import chromadb
import os
import json

# Load embedding model
model = SentenceTransformer("all-MiniLM-L6-v2")

# ChromaDB setup
client = chromadb.PersistentClient(path="vector_db")

collection = client.get_or_create_collection(
    name="knowledge_base"
)

DOCUMENTS_FOLDER = "documents"
PROCESSED_FILE = "processed_files.json"


def load_processed_files():
    if os.path.exists(PROCESSED_FILE):
        with open(PROCESSED_FILE, "r") as f:
            return json.load(f)
    return []


def save_processed_files(files):
    with open(PROCESSED_FILE, "w") as f:
        json.dump(files, f, indent=4)


processed_files = load_processed_files()

print("Checking for new files...\n")

for file_name in os.listdir(DOCUMENTS_FOLDER):

    if not file_name.endswith(".pptx"):
        continue

    if file_name in processed_files:
        print(f"Skipping: {file_name} (Already Processed)")
        continue

    file_path = os.path.join(DOCUMENTS_FOLDER, file_name)

    print(f"\nReading: {file_name}")

    presentation = Presentation(file_path)

    text = ""

    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"

    print("Text Extracted Successfully")

    splitter = RecursiveCharacterTextSplitter(
        chunk_size=500,
        chunk_overlap=100
    )

    chunks = splitter.split_text(text)

    filtered_chunks = []

    for chunk in chunks:
        chunk = chunk.strip()

        if len(chunk) > 100:
            filtered_chunks.append(chunk)

    chunks = filtered_chunks

    print("Total Chunks:", len(chunks))

    embeddings = model.encode(chunks)

    print("Embeddings Generated Successfully")

    for i, chunk in enumerate(chunks):

        collection.add(
            ids=[f"{file_name}_{i}"],
            documents=[chunk],
            embeddings=[embeddings[i].tolist()]
        )

    print("Stored in ChromaDB Successfully")

    processed_files.append(file_name)
    save_processed_files(processed_files)

print("\nDatabase Update Completed")